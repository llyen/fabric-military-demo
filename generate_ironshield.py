"""
Generate demo scenario documents for Microsoft Fabric in military context.
OPERATION IRONSHIELD - Battlefield Simulation Demo Scenario.
"""

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as I, Pt as P
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN
import os

OUT = os.path.dirname(os.path.abspath(__file__))

# ── Shared color palette ──
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
OLIVE = RGBColor(0x3D, 0x5A, 0x1E)
ACCENT = RGBColor(0xC7, 0x8C, 0x1E)  # gold/amber
RED_ALERT = RGBColor(0xCC, 0x22, 0x22)
PNAVY = PptxRGB(0x0B, 0x1D, 0x3A)
POLIVE = PptxRGB(0x3D, 0x5A, 0x1E)
PACCENT = PptxRGB(0xC7, 0x8C, 0x1E)
PWHITE = PptxRGB(0xFF, 0xFF, 0xFF)
PGRAY = PptxRGB(0xE8, 0xE8, 0xE8)

# ================================================================
# DOCX HELPERS
# ================================================================

def style_doc(doc):
    s = doc.styles['Normal']
    s.font.name = 'Calibri'
    s.font.size = Pt(10.5)
    s.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    s.paragraph_format.space_after = Pt(6)

def dh(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for r in h.runs:
        r.font.color.rgb = NAVY if level <= 2 else OLIVE
    return h

def db(doc, text, bold_pfx=None):
    p = doc.add_paragraph(style='List Bullet')
    if bold_pfx:
        r = p.add_run(bold_pfx)
        r.bold = True
        p.add_run(text)
    else:
        p.add_run(text)
    return p

def dtable(doc, headers, rows):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.style = 'Light Grid Accent 1'
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i]; c.text = h
        for p in c.paragraphs:
            for r in p.runs: r.bold = True; r.font.size = Pt(9)
    for ri, row in enumerate(rows):
        for ci, v in enumerate(row):
            c = t.rows[ri+1].cells[ci]; c.text = v
            for p in c.paragraphs:
                for r in p.runs: r.font.size = Pt(9)
    return t

def dcode(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.font.name = 'Consolas'
    r.font.size = Pt(8.5)
    r.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    return p

# ================================================================
# PPTX HELPERS
# ================================================================

def ptitle(prs, title, subtitle=""):
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = title
    for sh in sl.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = PNAVY
    if subtitle and sl.placeholders[1]:
        sl.placeholders[1].text = subtitle
    return sl

def pcontent(prs, title, bullets, layout=1):
    sl = prs.slides.add_slide(prs.slide_layouts[layout])
    sl.shapes.title.text = title
    for sh in sl.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size and r.font.size >= P(20):
                        r.font.color.rgb = PNAVY
    body = sl.placeholders[1]
    tf = body.text_frame; tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = b
            tf.paragraphs[0].font.size = P(14)
        else:
            p = tf.add_paragraph(); p.text = b; p.font.size = P(14); p.level = 0
    return sl

def p2col(prs, title, lt, lb, rt, rb):
    sl = prs.slides.add_slide(prs.slide_layouts[5])
    tx = sl.shapes.add_textbox(I(0.5), I(0.3), I(9), I(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = title; p.font.size = P(24); p.font.bold = True; p.font.color.rgb = PNAVY

    lbox = sl.shapes.add_textbox(I(0.5), I(1.3), I(4.3), I(5.5))
    ltf = lbox.text_frame; ltf.word_wrap = True
    lp = ltf.paragraphs[0]; lp.text = lt; lp.font.size = P(16); lp.font.bold = True; lp.font.color.rgb = POLIVE
    for b in lb:
        p = ltf.add_paragraph(); p.text = f"• {b}"; p.font.size = P(12); p.space_after = P(4)

    rbox = sl.shapes.add_textbox(I(5.2), I(1.3), I(4.3), I(5.5))
    rtf = rbox.text_frame; rtf.word_wrap = True
    rp = rtf.paragraphs[0]; rp.text = rt; rp.font.size = P(16); rp.font.bold = True; rp.font.color.rgb = POLIVE
    for b in rb:
        p = rtf.add_paragraph(); p.text = f"• {b}"; p.font.size = P(12); p.space_after = P(4)
    return sl


# ================================================================
# 1. MAIN DOCX – Scenario Document
# ================================================================

def create_scenario_docx():
    doc = Document()
    style_doc(doc)

    # ── Title page ──
    doc.add_paragraph(); doc.add_paragraph()
    t = doc.add_heading('OPERATION IRONSHIELD', level=0)
    for r in t.runs: r.font.color.rgb = NAVY; r.font.size = Pt(30)
    sub = doc.add_paragraph('Symulacja pola walki z wykorzystaniem\nMicrosoft Fabric Real-Time Intelligence')
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for r in sub.runs: r.font.size = Pt(14); r.font.color.rgb = OLIVE
    doc.add_paragraph()
    doc.add_paragraph('Scenariusz demonstracyjny: Fabric RTI • Data Agents • Operations Agents')
    doc.add_paragraph()
    p = doc.add_paragraph('DOKUMENT POUFNY')
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for r in p.runs: r.bold = True; r.font.color.rgb = RED_ALERT
    doc.add_page_break()

    # ── TOC ──
    dh(doc, 'Spis treści', 1)
    toc = [
        '1. Kontekst i cel demonstracji',
        '2. Scenariusz: OPERATION IRONSHIELD',
        '3. Architektura techniczna rozwiązania',
        '4. Faza 1 – Real-Time Intelligence: Świadomość sytuacyjna',
        '5. Faza 2 – Data Agents: Wywiad i analiza',
        '6. Faza 3 – Operations Agents: Autonomiczne reagowanie',
        '7. Scenariusz demonstracji krok po kroku',
        '8. Dane symulacyjne i generatory',
        '9. Wymagania techniczne do uruchomienia demo',
        '10. Podsumowanie wartości biznesowej',
    ]
    for t in toc: doc.add_paragraph(t)
    doc.add_page_break()

    # ═══════════════════════════════════════════════════
    # 1. KONTEKST
    # ═══════════════════════════════════════════════════
    dh(doc, '1. Kontekst i cel demonstracji', 1)
    doc.add_paragraph(
        'Współczesne pole walki generuje ogromne ilości danych w czasie rzeczywistym: '
        'telemetria z dronów, czujniki na pojazdach opancerzonych, dane radarowe, '
        'transmisje radiowe, dane GPS żołnierzy, obrazy satelitarne i dane meteorologiczne. '
        'Zdolność do natychmiastowego przetwarzania, korelowania i podejmowania decyzji '
        'na podstawie tych danych stanowi kluczową przewagę taktyczną.'
    )
    doc.add_paragraph(
        'Celem demonstracji jest pokazanie, jak Microsoft Fabric – jako zunifikowana '
        'platforma danych i AI – może pełnić rolę rdzenia systemu C4ISR '
        '(Command, Control, Communications, Computers, Intelligence, Surveillance, Reconnaissance), '
        'dostarczając:'
    )
    db(doc, ' – ciągłą świadomość sytuacyjną (Real-Time Intelligence)', 'Obraz pola walki')
    db(doc, ' – konwersacyjny dostęp do danych wywiadu (Data Agents)', 'Wsparcie decyzyjne')
    db(doc, ' – automatyczną reakcję na zdarzenia krytyczne (Operations Agents)', 'Autonomiczne reagowanie')

    dh(doc, '1.1 Dlaczego Fabric?', 2)
    dtable(doc,
        ['Zdolność C4ISR', 'Komponent Fabric', 'Korzyść'],
        [
            ['Zbieranie danych sensorowych', 'Eventstream / Real-Time Hub',
             'Ingestion milionów zdarzeń/s z heterogenicznych źródeł'],
            ['Przechowywanie i analiza', 'Eventhouse (KQL Database)',
             'Sub-sekundowe zapytania na miliardach rekordów'],
            ['Wizualizacja sytuacyjna', 'Real-Time Dashboard / Power BI',
             'Mapy operacyjne, KPI, alerty w czasie rzeczywistym'],
            ['Wykrywanie anomalii', 'AI-powered Anomaly Detection',
             'Automatyczne wykrycie nietypowych wzorców ruchu wroga'],
            ['Wywiad konwersacyjny', 'Data Agents',
             'Dowódca pyta w języku naturalnym, agent odpowiada danymi'],
            ['Reakcja autonomiczna', 'Operations Agents + Activator',
             'Cykl observe→analyze→decide→act z human-in-the-loop'],
            ['Governance i bezpieczeństwo', 'Unity Catalog / RBAC',
             'Klasyfikacja danych, kontrola dostępu, audyt'],
        ]
    )

    # ═══════════════════════════════════════════════════
    # 2. SCENARIUSZ
    # ═══════════════════════════════════════════════════
    dh(doc, '2. Scenariusz: OPERATION IRONSHIELD', 1)
    doc.add_paragraph(
        'Scenariusz symuluje operację obronną brygady zmechanizowanej na terenie '
        'zurbanizowanym. Siły sojusznicze (BLUE FORCE) bronią strategicznego węzła '
        'komunikacyjnego przed atakiem sił przeciwnika (RED FORCE).'
    )

    dh(doc, '2.1 Siły i zasoby BLUE FORCE', 2)
    dtable(doc,
        ['Kategoria', 'Zasób', 'Ilość', 'Sensory / dane'],
        [
            ['Pojazdy opancerzone', 'BWP Borsuk / Rosomak', '24', 'GPS, telemetria silnika, status amunicji, kamera termowizyjna'],
            ['Drony rozpoznawcze', 'FlyEye / Warmate', '12', 'Video HD, IR, GPS, prędkość, wysokość, bateria'],
            ['Żołnierze piechoty', 'Sekcje piechoty', '120', 'GPS (BMS), tętno, status zdrowia (wearable)'],
            ['Radary', 'Radar artyleryjski', '4', 'Wykrycia obiektów, tracking, klasyfikacja'],
            ['Artyleria', 'Haubice Krab', '8', 'Status gotowości, amunicja, pozycja'],
            ['Stacje meteorologiczne', 'Automatyczne stacje', '6', 'Temperatura, wiatr, widoczność, opady'],
            ['Logistyka', 'Pojazdy zaopatrzeniowe', '10', 'GPS, ładunek, trasa, ETA'],
        ]
    )

    dh(doc, '2.2 Siły RED FORCE (symulowane)', 2)
    doc.add_paragraph(
        'Przeciwnik jest symulowany przez generator danych, który emituje zdarzenia: '
        'wykrycia radarowe, obserwacje z dronów, przechwycone transmisje radiowe (SIGINT). '
        'Zachowanie wroga jest parametryzowane i może obejmować:'
    )
    db(doc, 'Marsz kolumny pancernej w kierunku węzła komunikacyjnego')
    db(doc, 'Rekonesans dronami w strefie BLUE FORCE')
    db(doc, 'Atak artyleryjski na pozycje sojusznicze')
    db(doc, 'Próba infiltracji grupy dywersyjnej')
    db(doc, 'Zakłócanie łączności (EW – Electronic Warfare)')

    dh(doc, '2.3 Oś czasu scenariusza', 2)
    dtable(doc,
        ['Czas (T+)', 'Wydarzenie', 'Demonstrowana zdolność Fabric'],
        [
            ['T+0:00', 'Inicjalizacja: deploy sensorów, start streamów', 'Eventstream – ingestion z wielu źródeł'],
            ['T+2:00', 'Normalna aktywność, patrol, dane bazowe', 'RTI Dashboard – mapa operacyjna live'],
            ['T+5:00', 'Radar wykrywa niezidentyfikowane obiekty w sektorze Bravo', 'Anomaly Detection – alert automatyczny'],
            ['T+6:00', 'Dowódca pyta: "Ile kontaktów? Klasyfikacja? Kierunek?"', 'Data Agent – zapytanie NL'],
            ['T+7:00', 'Data Agent koreluje dane radar + dron + SIGINT', 'Data Agent – multi-source reasoning'],
            ['T+8:00', 'Identyfikacja: kolumna pancerna wroga, 14 pojazdów', 'RTI Dashboard update + alert'],
            ['T+9:00', 'Operations Agent rekomenduje: alarmowanie artylerii', 'Operations Agent – automated recommendation'],
            ['T+10:00', 'Human-in-the-loop: dowódca zatwierdza w Teams', 'Operations Agent – approval workflow'],
            ['T+11:00', 'Automatyczne przekazanie współrzędnych do baterii Krab', 'Operations Agent → Power Automate action'],
            ['T+12:00', 'Monitoring efektów: dron weryfikuje trafienia', 'RTI – Battle Damage Assessment'],
            ['T+15:00', 'Data Agent: "Podsumuj sytuację w sektorze Bravo"', 'Data Agent – summary & report generation'],
            ['T+18:00', 'Operations Agent wykrywa niski stan amunicji 3. kompanii', 'Operations Agent – logistics optimization'],
            ['T+20:00', 'Automatyczna optymalizacja trasy konwoju zaopatrzeniowego', 'Operations Agent → routing action'],
        ]
    )

    # ═══════════════════════════════════════════════════
    # 3. ARCHITEKTURA
    # ═══════════════════════════════════════════════════
    dh(doc, '3. Architektura techniczna rozwiązania', 1)
    doc.add_paragraph(
        'Architektura opiera się na trzech warstwach Fabric zintegrowanych w jeden ekosystem:'
    )
    dcode(doc,
        '┌─────────────────────────────────────────────────────────────────────┐\n'
        '│                        ŹRÓDŁA DANYCH (SENSORY)                     │\n'
        '│  Drony   Radary   Pojazdy   Żołnierze   Meteo   SIGINT   Logistyka│\n'
        '└──────────┬──────────┬───────────┬──────────┬──────────┬────────────┘\n'
        '           │ MQTT/Kafka│  IoT Hub  │   REST   │ Event Hub│\n'
        '           ▼          ▼           ▼          ▼          ▼\n'
        '┌─────────────────────────────────────────────────────────────────────┐\n'
        '│                    FABRIC REAL-TIME INTELLIGENCE                    │\n'
        '│  ┌──────────────┐  ┌────────────────┐  ┌───────────────────────┐   │\n'
        '│  │  Eventstream  │→│   Eventhouse    │→│  Real-Time Dashboard  │   │\n'
        '│  │  (ingestion)  │ │  (KQL Database) │ │  (mapa operacyjna)    │   │\n'
        '│  └──────────────┘  └───────┬────────┘  └───────────────────────┘   │\n'
        '│                            │                                        │\n'
        '│  ┌─────────────────────────┼─────────────────────────────────────┐  │\n'
        '│  │         DATA AGENTS     │      OPERATIONS AGENTS              │  │\n'
        '│  │  ┌─────────────────┐    │   ┌──────────────────────────────┐  │  │\n'
        '│  │  │ Wywiad Agent    │    │   │ Threat Response Agent        │  │  │\n'
        '│  │  │ (NL → KQL/SQL)  │    │   │ (observe→analyze→decide→act) │  │  │\n'
        '│  │  ├─────────────────┤    │   ├──────────────────────────────┤  │  │\n'
        '│  │  │ Logistics Agent │    │   │ Logistics Ops Agent          │  │  │\n'
        '│  │  │ (zapasy, trasy) │    │   │ (resupply optimization)      │  │  │\n'
        '│  │  └─────────────────┘    │   └──────────────┬───────────────┘  │  │\n'
        '│  └─────────────────────────┼──────────────────┼──────────────────┘  │\n'
        '│                            │                  │                      │\n'
        '│  ┌─────────────────────────┼──────────────────┼──────────────────┐  │\n'
        '│  │                 Unity Catalog (Governance)                     │  │\n'
        '│  │           Klasyfikacja • RBAC • Audyt • Lineage               │  │\n'
        '│  └───────────────────────────────────────────────────────────────┘  │\n'
        '└──────────────────────────────────┬──────────────────────────────────┘\n'
        '                                   │\n'
        '                    ┌───────────────┼───────────────┐\n'
        '                    │       Power Automate          │\n'
        '                    │  Teams alerts • Actions • C2  │\n'
        '                    └───────────────────────────────┘'
    )

    dh(doc, '3.1 Komponenty Fabric w scenariuszu', 2)

    dh(doc, 'Eventstream', 3)
    doc.add_paragraph(
        'Eventstream przyjmuje strumienie danych z symulowanych sensorów. '
        'Każdy typ sensora ma dedykowany strumień z transformacjami inline:'
    )
    dtable(doc,
        ['Strumień', 'Źródło', 'Format', 'Transformacje inline'],
        [
            ['drone-telemetry', 'Event Hub (symulator)', 'JSON', 'Parsowanie GPS, normalizacja altitude'],
            ['vehicle-status', 'Event Hub (symulator)', 'JSON', 'Enrichment: unit_id → unit_name'],
            ['radar-detections', 'Event Hub (symulator)', 'JSON', 'Filtrowanie szumu (confidence > 0.7)'],
            ['soldier-health', 'Event Hub (symulator)', 'JSON', 'Anonymization, alert threshold (HR>180)'],
            ['weather-data', 'Event Hub (symulator)', 'JSON', 'Agregacja 1-min okno'],
            ['sigint-intercepts', 'Event Hub (symulator)', 'JSON', 'Classification tag'],
            ['ammo-logistics', 'Event Hub (symulator)', 'JSON', 'Enrichment: depot_location'],
        ]
    )

    dh(doc, 'Eventhouse (KQL Database)', 3)
    doc.add_paragraph('KQL Database przechowuje wszystkie strumienie i udostępnia sub-sekundowe zapytania:')
    dcode(doc,
        '// Przykład: Wykrywanie koncentracji sił wroga w ostatnich 10 minutach\n'
        'RadarDetections\n'
        '| where Timestamp > ago(10m)\n'
        '| where Classification == "hostile"\n'
        '| summarize ContactCount = count(), \n'
        '           AvgSpeed = avg(Speed_kmh),\n'
        '           Centroid = avgif(pack("lat", Latitude, "lon", Longitude), true)\n'
        '  by Sector, bin(Timestamp, 1m)\n'
        '| where ContactCount > 5\n'
        '| order by ContactCount desc'
    )
    dcode(doc,
        '// Przykład: Status zdrowotny żołnierzy z alertami\n'
        'SoldierHealth\n'
        '| where Timestamp > ago(5m)\n'
        '| summarize LatestHR = arg_max(Timestamp, HeartRate),\n'
        '            LatestTemp = arg_max(Timestamp, BodyTemp)\n'
        '  by SoldierId, UnitName\n'
        '| where LatestHR > 180 or LatestTemp > 39.5\n'
        '| project SoldierId, UnitName, LatestHR, LatestTemp,\n'
        '          AlertLevel = case(LatestHR > 200, "CRITICAL", "WARNING")'
    )
    dcode(doc,
        '// Przykład: Stan amunicji floty pojazdów\n'
        'VehicleStatus\n'
        '| where Timestamp > ago(2m)\n'
        '| summarize arg_max(Timestamp, *) by VehicleId\n'
        '| project VehicleId, UnitName, AmmoPercent, FuelPercent,\n'
        '          CombatReady = AmmoPercent > 30 and FuelPercent > 20\n'
        '| summarize ReadyCount = countif(CombatReady),\n'
        '            TotalCount = count()\n'
        '  by UnitName'
    )

    dh(doc, 'Real-Time Dashboard', 3)
    doc.add_paragraph('Dashboard operacyjny zawiera następujące panele (tiles):')
    dtable(doc,
        ['Panel', 'Typ wizualizacji', 'Źródło danych', 'Odświeżanie'],
        [
            ['Mapa operacyjna', 'Mapa geospatialna', 'Wszystkie strumienie GPS', 'Ciągłe (1s)'],
            ['Status sił BLUE', 'Scorecard / KPI', 'VehicleStatus + SoldierHealth', '5s'],
            ['Wykrycia RED FORCE', 'Scatter plot na mapie', 'RadarDetections + DroneObservations', '2s'],
            ['Alert feed', 'Lista zdarzeń', 'Anomaly Detection output', 'Ciągłe'],
            ['Pogoda operacyjna', 'Gauge + map overlay', 'WeatherData', '60s'],
            ['Logistyka', 'Bar chart', 'AmmoLogistics', '10s'],
            ['Oś czasu incydentów', 'Timeline', 'Wszystkie alerty', 'Ciągłe'],
        ]
    )

    # ═══════════════════════════════════════════════════
    # 4. FAZA 1 – RTI
    # ═══════════════════════════════════════════════════
    dh(doc, '4. Faza 1 – Real-Time Intelligence: Świadomość sytuacyjna', 1)
    doc.add_paragraph(
        'W tej fazie demonstracji pokazujemy, jak Fabric RTI zbiera, przetwarza i wizualizuje '
        'dane z pola walki w czasie rzeczywistym, tworząc Common Operational Picture (COP).'
    )

    dh(doc, '4.1 Co pokazujemy', 2)
    db(doc, 'Ingestion z 7 strumieni danych jednocześnie (Eventstream)')
    db(doc, 'Live mapa operacyjna z pozycjami BLUE i RED FORCE')
    db(doc, 'Automatyczne wykrycie anomalii: nagły wzrost kontaktów radarowych w sektorze Bravo')
    db(doc, 'Alert w czasie rzeczywistym na dashboardzie i w Teams')
    db(doc, 'Drill-down: kliknięcie na alert → szczegóły KQL')

    dh(doc, '4.2 Anomaly Detection – konfiguracja', 2)
    doc.add_paragraph(
        'Fabric RTI oferuje wbudowaną detekcję anomalii na strumieniach danych. '
        'W scenariuszu konfigurujemy ją dla strumienia RadarDetections:'
    )
    db(doc, ' – liczba wykryć hostile w 2-minutowym oknie', 'Metryka:')
    db(doc, ' – AnomalyScore > 0.8 (odchylenie od baseline)', 'Próg alertu:')
    db(doc, ' – Activator → Teams notification do dowódcy + Operations Agent trigger', 'Akcja:')

    dh(doc, '4.3 Narracja dla demonstracji', 2)
    doc.add_paragraph(
        '„Widzimy na żywo mapę operacyjną. Niebieskie ikony to nasze siły – pojazdy, drony, '
        'żołnierze. Każda ikona aktualizuje się co sekundę na podstawie danych GPS. '
        'W prawym panelu widzimy status sił: 22 z 24 pojazdów combat-ready, 118 żołnierzy z '
        'normalnymi parametrami zdrowotnymi.'
    )
    doc.add_paragraph(
        'Nagle w sektorze Bravo pojawia się anomalia – system wykrywa gwałtowny wzrost '
        'kontaktów radarowych. Alert pojawia się automatycznie na dashboardzie i jako '
        'notyfikacja w Microsoft Teams. To przejście do Fazy 2."'
    )

    # ═══════════════════════════════════════════════════
    # 5. FAZA 2 – DATA AGENTS
    # ═══════════════════════════════════════════════════
    dh(doc, '5. Faza 2 – Data Agents: Wywiad i analiza', 1)
    doc.add_paragraph(
        'Data Agents w Microsoft Fabric umożliwiają dowódcy zadawanie pytań w języku naturalnym '
        'i otrzymywanie odpowiedzi opartych na aktualnych danych z Eventhouse. '
        'Agent tłumaczy zapytanie na KQL/SQL, wykonuje je i prezentuje wynik w zrozumiałej formie.'
    )

    dh(doc, '5.1 Konfiguracja Data Agents', 2)
    dtable(doc,
        ['Agent', 'Nazwa', 'Źródła danych', 'Rola'],
        [
            ['Wywiad Agent', 'IRONSHIELD Intel', 'RadarDetections, DroneObservations, SigintIntercepts',
             'Analiza zagrożeń, korelacja danych wywiadowczych'],
            ['Logistics Agent', 'IRONSHIELD Supply', 'AmmoLogistics, VehicleStatus, WeatherData',
             'Status zaopatrzenia, optymalizacja tras dostaw'],
            ['Health Agent', 'IRONSHIELD Medic', 'SoldierHealth, VehicleStatus',
             'Monitoring stanu zdrowia, rekomendacje MEDEVAC'],
        ]
    )

    dh(doc, '5.2 Przykładowe interakcje', 2)
    doc.add_paragraph('Poniżej przykładowe zapytania dowódcy i odpowiedzi agentów:')

    doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run('DOWÓDCA: '); r.bold = True; r.font.color.rgb = RGBColor(0x00, 0x50, 0x99)
    p.add_run('"Ile kontaktów wrogich wykryto w sektorze Bravo w ostatnich 10 minutach? '
              'Jaki jest ich kierunek i prędkość?"')

    p = doc.add_paragraph()
    r = p.add_run('INTEL AGENT: '); r.bold = True; r.font.color.rgb = OLIVE
    p.add_run('"W sektorze Bravo wykryto 14 kontaktów klasyfikowanych jako hostile. '
              'Średnia prędkość: 35 km/h, kierunek: północny-zachód (azymut 315°). '
              'Korelacja z danymi SIGINT wskazuje na transmisje radiowe odpowiadające '
              'dowodzeniu kolumną pancerną. Dystans do naszych pozycji: ~12 km. '
              'Szacowany czas dotarcia: 20 minut przy obecnej prędkości."')

    doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run('DOWÓDCA: '); r.bold = True; r.font.color.rgb = RGBColor(0x00, 0x50, 0x99)
    p.add_run('"Jakie siły mamy w zasięgu reagowania? Stan amunicji?"')

    p = doc.add_paragraph()
    r = p.add_run('SUPPLY AGENT: '); r.bold = True; r.font.color.rgb = OLIVE
    p.add_run('"W zasięgu 15 km od sektora Bravo: 2. kompania zmechanizowana (8 BWP, combat-ready), '
              'bateria Krab (4 haubice, amunicja 78%). 3. kompania ma stan amunicji 25% – poniżej progu '
              'operacyjnego. Konwój zaopatrzeniowy KZ-04 jest w trasie, ETA 45 minut."')

    doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run('DOWÓDCA: '); r.bold = True; r.font.color.rgb = RGBColor(0x00, 0x50, 0x99)
    p.add_run('"Czy warunki pogodowe pozwalają na wsparcie lotnicze w sektorze Bravo?"')

    p = doc.add_paragraph()
    r = p.add_run('INTEL AGENT: '); r.bold = True; r.font.color.rgb = OLIVE
    p.add_run('"Aktualne warunki w sektorze Bravo: zachmurzenie 6/8, podstawa chmur 800m, '
              'widoczność 5km, wiatr SW 15 km/h. Prognoza na najbliższą godzinę: pogorszenie widoczności '
              'do 2km. Operacje dronów: możliwe. Wsparcie lotnicze samolotami: ograniczone ze względu na '
              'niską podstawę chmur."')

    dh(doc, '5.3 Mechanizm działania (za kulisami)', 2)
    doc.add_paragraph(
        'Data Agent tłumaczy pytanie na zapytanie KQL, wykonuje je na Eventhouse, '
        'i formatuje odpowiedź w języku naturalnym. Przykład tłumaczenia:'
    )
    dcode(doc,
        '// Zapytanie wygenerowane przez Intel Agent:\n'
        'let timeWindow = ago(10m);\n'
        'let sectorFilter = "Bravo";\n'
        'RadarDetections\n'
        '| where Timestamp > timeWindow and Sector == sectorFilter\n'
        '| where Classification == "hostile"\n'
        '| summarize \n'
        '    ContactCount = dcount(TrackId),\n'
        '    AvgSpeed = avg(Speed_kmh),\n'
        '    AvgHeading = avg(Heading_deg),\n'
        '    MinDistance = min(DistanceToBlue_km)\n'
        '| extend ETA_min = MinDistance / AvgSpeed * 60'
    )

    dh(doc, '5.4 Narracja', 2)
    doc.add_paragraph(
        '„Dowódca otrzymał alert o anomalii w sektorze Bravo. Zamiast ręcznie przeszukiwać '
        'tabele i wykresy, po prostu pyta agenta w języku naturalnym. Agent w tle wykonuje '
        'zaawansowane zapytania KQL, koreluje dane z radaru, dronów i SIGINT, i dostarcza '
        'syntetyczną odpowiedź z rekomendacją. To skraca pętlę OODA z minut do sekund."'
    )

    # ═══════════════════════════════════════════════════
    # 6. FAZA 3 – OPERATIONS AGENTS
    # ═══════════════════════════════════════════════════
    dh(doc, '6. Faza 3 – Operations Agents: Autonomiczne reagowanie', 1)
    doc.add_paragraph(
        'Operations Agents w Fabric realizują pełny cykl autonomicznego reagowania: '
        'observe → analyze → decide → act. Agent monitoruje dane w Eventhouse, '
        'rozpoznaje sytuacje wymagające akcji, formułuje rekomendację i – po zatwierdzeniu '
        'przez człowieka – uruchamia workflow Power Automate.'
    )

    dh(doc, '6.1 Konfiguracja Operations Agents', 2)
    dtable(doc,
        ['Agent', 'Cel biznesowy', 'Źródło danych', 'Akcja rekomendowana', 'Kanał zatwierdzenia'],
        [
            ['Threat Response', 'Neutralizacja zagrożenia',
             'RadarDetections + DroneObservations',
             'Alert artylerii, przekazanie współrzędnych',
             'Teams – Adaptive Card'],
            ['Logistics Optimizer', 'Utrzymanie zdolności bojowej',
             'AmmoLogistics + VehicleStatus',
             'Rerouting konwoju, priorytetyzacja dostaw',
             'Teams – Adaptive Card'],
            ['MEDEVAC Coordinator', 'Ochrona życia żołnierzy',
             'SoldierHealth',
             'Wezwanie MEDEVAC, wskazanie LZ',
             'Teams – natychmiastowy alert'],
            ['EW Alert', 'Ochrona przed zakłócaniem',
             'SigintIntercepts + Comms quality',
             'Przełączenie na częstotliwości zapasowe',
             'Automatyczny (pre-approved)'],
        ]
    )

    dh(doc, '6.2 Przykład: Threat Response Agent – pełny cykl', 2)

    dh(doc, 'OBSERVE', 3)
    doc.add_paragraph(
        'Agent ciągle monitoruje zapytanie KQL na Eventhouse, sprawdzając liczbę '
        'kontaktów hostile z dużą prędkością zbliżającą się do pozycji BLUE FORCE.'
    )
    dcode(doc,
        '// Monitoring query (co 30s)\n'
        'RadarDetections\n'
        '| where Timestamp > ago(5m) and Classification == "hostile"\n'
        '| summarize Count = dcount(TrackId), \n'
        '            AvgDist = avg(DistanceToBlue_km)\n'
        '  by Sector\n'
        '| where Count >= 10 and AvgDist < 15'
    )

    dh(doc, 'ANALYZE', 3)
    doc.add_paragraph(
        'Gdy warunek się spełnia, agent rozpoznaje sytuację jako potencjalny atak pancerny. '
        'Analizuje kontekst: dostępne środki ogniowe, warunki pogodowe, pozycje sił własnych.'
    )

    dh(doc, 'DECIDE', 3)
    doc.add_paragraph(
        'Agent formułuje rekomendację i wysyła ją jako Adaptive Card w Microsoft Teams:'
    )
    doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run('⚠️ THREAT RESPONSE RECOMMENDATION\n\n'); r.bold = True; r.font.color.rgb = RED_ALERT
    p.add_run(
        'Wykryto kolumnę pancerną wroga (14 pojazdów) w sektorze Bravo.\n'
        'Azymut: 315° | Prędkość: 35 km/h | ETA do naszych pozycji: 20 min\n\n'
        'REKOMENDACJA: Ogień zaporowy artylerii (Bateria Krab, poz. K-7)\n'
        'Współrzędne celu: 51.234°N, 21.456°E\n'
        'Szacowane zużycie amunicji: 12 pocisków\n\n'
        '[✅ ZATWIERDŹ]    [❌ ODRZUĆ]    [✏️ MODYFIKUJ]'
    )

    dh(doc, 'ACT (po zatwierdzeniu)', 3)
    doc.add_paragraph(
        'Po kliknięciu „ZATWIERDŹ" przez dowódcę w Teams, Power Automate workflow:'
    )
    db(doc, 'Przekazuje współrzędne celu do systemu ogniowego (symulowane API)')
    db(doc, 'Aktualizuje status artylerii w Eventhouse')
    db(doc, 'Wysyła dron rozpoznawczy do weryfikacji efektów (BDA)')
    db(doc, 'Loguje decyzję w Unity Catalog (audyt)')

    dh(doc, '6.3 Przykład: Logistics Optimizer Agent', 2)
    doc.add_paragraph(
        'Agent monitoruje stan amunicji i paliwa wszystkich jednostek. Gdy 3. kompania '
        'spada poniżej 30% amunicji, agent:'
    )
    db(doc, 'Identyfikuje najbliższy konwój zaopatrzeniowy (KZ-04)')
    db(doc, 'Optymalizuje trasę uwzględniając pozycje wroga i warunki drogowe')
    db(doc, 'Proponuje zmianę priorytetu ładunku (amunicja > paliwo)')
    db(doc, 'Wysyła rekomendację do oficera logistyki w Teams')
    db(doc, 'Po zatwierdzeniu – aktualizuje routing konwoju w systemie nawigacji')

    dh(doc, '6.4 Narracja', 2)
    doc.add_paragraph(
        '„Operations Agent działa jak inteligentny oficer operacyjny, który nigdy nie śpi. '
        'Ciągle monitoruje setki parametrów, rozpoznaje sytuacje wymagające reakcji i '
        'proponuje konkretne akcje. Ale kluczowa decyzja zawsze pozostaje w rękach człowieka – '
        'dowódca zatwierdza lub modyfikuje rekomendację jednym kliknięciem w Teams. '
        'To jest human-in-the-loop AI na polu walki."'
    )

    # ═══════════════════════════════════════════════════
    # 7. KROK PO KROKU
    # ═══════════════════════════════════════════════════
    dh(doc, '7. Scenariusz demonstracji krok po kroku', 1)
    doc.add_paragraph(
        'Poniżej szczegółowy scenariusz do przeprowadzenia demo (~20-25 minut):'
    )

    steps = [
        ('Otwarcie (2 min)',
         'Kontekst: „Wyobraźmy sobie, że jesteśmy w centrum dowodzenia brygady zmechanizowanej. '
         'Przed nami operacja IRONSHIELD – obrona strategicznego węzła komunikacyjnego."'),
        ('Start symulatora (1 min)',
         'Uruchomienie generatora danych (skrypt Python/notebook Fabric). '
         'Pokazanie Eventstream: 7 strumieni, dane wpływają w czasie rzeczywistym.'),
        ('Real-Time Dashboard (3 min)',
         'Otwarcie dashboardu operacyjnego. Omówienie paneli: mapa, KPI, alert feed. '
         'Pokazanie live aktualizacji – ikony poruszają się na mapie.'),
        ('Anomalia (2 min)',
         'Symulator wprowadza zdarzenie: kolumna pancerna wroga w sektorze Bravo. '
         'Anomaly Detection wykrywa i generuje alert na dashboardzie + Teams.'),
        ('Data Agent – wywiad (4 min)',
         'Otwarcie Data Agent (IRONSHIELD Intel). Zadanie 2-3 pytań w języku naturalnym. '
         'Pokazanie jak agent tłumaczy na KQL i dostarcza syntetyczną odpowiedź.'),
        ('Data Agent – logistyka (2 min)',
         'Przełączenie na Supply Agent. Pytanie o stan amunicji i dostępne środki ogniowe.'),
        ('Operations Agent – Threat Response (4 min)',
         'Pokazanie Adaptive Card w Teams z rekomendacją. '
         'Omówienie cyklu observe→analyze→decide→act. '
         'Zatwierdzenie akcji przez dowódcę. Pokazanie uruchomionego workflow.'),
        ('Operations Agent – Logistics (3 min)',
         'Pokazanie automatycznej detekcji niskiego stanu amunicji. '
         'Rekomendacja reroutingu konwoju. Zatwierdzenie w Teams.'),
        ('Battle Damage Assessment (2 min)',
         'Powrót do dashboardu: dane z drona potwierdzają efekt. '
         'Data Agent podsumowuje sytuację.'),
        ('Podsumowanie wartości (2 min)',
         'Kluczowe przekazy: skrócenie OODA loop, human-in-the-loop, '
         'jednolita platforma, governance i audyt.'),
    ]

    for i, (title, desc) in enumerate(steps, 1):
        dh(doc, f'Krok {i}: {title}', 3)
        doc.add_paragraph(desc)

    # ═══════════════════════════════════════════════════
    # 8. DANE SYMULACYJNE
    # ═══════════════════════════════════════════════════
    dh(doc, '8. Dane symulacyjne i generatory', 1)
    doc.add_paragraph(
        'Do demonstracji potrzebny jest generator danych symulujący strumienie z sensorów. '
        'Poniżej schematy JSON dla głównych strumieni.'
    )

    dh(doc, '8.1 Schema: RadarDetection', 2)
    dcode(doc,
        '{\n'
        '  "EventId": "RD-2026-03-26-00142",\n'
        '  "Timestamp": "2026-03-26T10:05:23.456Z",\n'
        '  "Sector": "Bravo",\n'
        '  "TrackId": "TRK-0047",\n'
        '  "Classification": "hostile",      // hostile | friendly | unknown\n'
        '  "ObjectType": "armored_vehicle",   // armored_vehicle | infantry | drone | artillery\n'
        '  "Latitude": 51.2345,\n'
        '  "Longitude": 21.4567,\n'
        '  "Speed_kmh": 35.2,\n'
        '  "Heading_deg": 315,\n'
        '  "DistanceToBlue_km": 12.4,\n'
        '  "Confidence": 0.92,\n'
        '  "RadarId": "RAD-02"\n'
        '}'
    )

    dh(doc, '8.2 Schema: VehicleStatus', 2)
    dcode(doc,
        '{\n'
        '  "EventId": "VS-2026-03-26-08821",\n'
        '  "Timestamp": "2026-03-26T10:05:24.123Z",\n'
        '  "VehicleId": "BWP-2K-07",\n'
        '  "VehicleType": "BWP_Borsuk",\n'
        '  "UnitName": "2. kompania zmech.",\n'
        '  "Latitude": 51.2401,\n'
        '  "Longitude": 21.4100,\n'
        '  "Speed_kmh": 0,\n'
        '  "Heading_deg": 270,\n'
        '  "EngineStatus": "running",\n'
        '  "FuelPercent": 72,\n'
        '  "AmmoPercent": 85,\n'
        '  "CrewCount": 3,\n'
        '  "CombatReady": true\n'
        '}'
    )

    dh(doc, '8.3 Schema: SoldierHealth', 2)
    dcode(doc,
        '{\n'
        '  "EventId": "SH-2026-03-26-44210",\n'
        '  "Timestamp": "2026-03-26T10:05:25.789Z",\n'
        '  "SoldierId": "SOL-3K-042",\n'
        '  "UnitName": "3. kompania zmech.",\n'
        '  "Latitude": 51.2380,\n'
        '  "Longitude": 21.4250,\n'
        '  "HeartRate": 95,\n'
        '  "BodyTemp": 36.8,\n'
        '  "BloodOxygen": 98,\n'
        '  "StressLevel": "normal",     // normal | elevated | critical\n'
        '  "MovementStatus": "stationary" // stationary | walking | running | prone\n'
        '}'
    )

    dh(doc, '8.4 Schema: DroneObservation', 2)
    dcode(doc,
        '{\n'
        '  "EventId": "DO-2026-03-26-01523",\n'
        '  "Timestamp": "2026-03-26T10:05:26.456Z",\n'
        '  "DroneId": "FE-04",\n'
        '  "DroneType": "FlyEye",\n'
        '  "Latitude": 51.2350,\n'
        '  "Longitude": 21.4560,\n'
        '  "Altitude_m": 250,\n'
        '  "BatteryPercent": 68,\n'
        '  "ObservationType": "visual_contact",\n'
        '  "TargetClassification": "armored_column",\n'
        '  "TargetCount": 14,\n'
        '  "TargetLatitude": 51.2345,\n'
        '  "TargetLongitude": 21.4567,\n'
        '  "ImageUrl": "blob://observations/DO-01523.jpg",\n'
        '  "Confidence": 0.95\n'
        '}'
    )

    # ═══════════════════════════════════════════════════
    # 9. WYMAGANIA TECHNICZNE
    # ═══════════════════════════════════════════════════
    dh(doc, '9. Wymagania techniczne do uruchomienia demo', 1)
    dtable(doc,
        ['Komponent', 'Wymaganie', 'Uwagi'],
        [
            ['Microsoft Fabric', 'Capacity F2 lub wyższa', 'Data Agents wymagają F2+'],
            ['Workspace', 'Dedykowany workspace „IRONSHIELD"', 'Separation of concerns'],
            ['Eventstream', '7 strumieni skonfigurowanych', 'Źródło: Event Hub (symulator)'],
            ['Event Hub', 'Azure Event Hub namespace', 'Do symulacji – może być Basic tier'],
            ['Eventhouse', '1 Eventhouse z KQL Database', 'Tabele per strumień'],
            ['Real-Time Dashboard', '1 dashboard z 7 panelami', 'Konfiguracja opisana w sekcji 3'],
            ['Data Agents', '3 agenty skonfigurowane', 'Intel, Supply, Medic'],
            ['Operations Agents', '4 agenty skonfigurowane', 'Threat, Logistics, MEDEVAC, EW'],
            ['Power Automate', '2 flows (threat response, logistics)', 'Triggered by Operations Agent'],
            ['Microsoft Teams', 'Kanał „IRONSHIELD-OPS"', 'Adaptive Cards'],
            ['Generator danych', 'Python notebook w Fabric', 'Parametryzowane scenariusze'],
            ['Tenant settings', 'Copilot/AI enabled', 'Admin portal'],
        ]
    )

    # ═══════════════════════════════════════════════════
    # 10. PODSUMOWANIE WARTOŚCI
    # ═══════════════════════════════════════════════════
    dh(doc, '10. Podsumowanie wartości biznesowej', 1)
    doc.add_paragraph(
        'Demonstracja OPERATION IRONSHIELD pokazuje, jak Microsoft Fabric transformuje '
        'surowe dane sensorowe w przewagę informacyjną i decyzyjną na polu walki.'
    )

    dh(doc, 'Kluczowe przekazy', 2)
    dtable(doc,
        ['Przekaz', 'Fabric = odpowiedź'],
        [
            ['Szybkość decyzji (OODA loop)',
             'Od detekcji do rekomendacji w sekundach, nie minutach. '
             'Data Agent + Operations Agent skracają pętlę decyzyjną.'],
            ['Jedna platforma, wiele zdolności',
             'RTI + AI Agents + Dashboards + Automation w jednym ekosystemie. '
             'Brak konieczności integracji wielu narzędzi.'],
            ['Human-in-the-loop',
             'AI rekomenduje, człowiek decyduje. Pełna kontrola dowódcy '
             'nad każdą krytyczną akcją – zatwierdzenie w Teams.'],
            ['Governance i audyt',
             'Unity Catalog: kto, kiedy, na jakiej podstawie podjął decyzję. '
             'Pełna traceability dla AAR (After Action Review).'],
            ['Skalowalność',
             'Od brygady do dywizji – ta sama architektura skaluje się '
             'na większą liczbę sensorów i jednostek.'],
            ['Interoperacyjność',
             'Open standards, REST API, Power Platform. '
             'Integracja z istniejącymi systemami C2 (BMS, JC3IEDM).'],
        ]
    )

    path = os.path.join(OUT, 'IRONSHIELD_Scenario_Demo.docx')
    doc.save(path)
    print(f"Created: {path}")
    return path


# ================================================================
# 2. PPTX – Presentation
# ================================================================

def create_scenario_pptx():
    prs = Presentation()
    prs.slide_width = I(13.333)
    prs.slide_height = I(7.5)

    # -- Slide 1: Title --
    ptitle(prs, 'OPERATION IRONSHIELD', 
           'Symulacja pola walki z Microsoft Fabric\nReal-Time Intelligence • Data Agents • Operations Agents')

    # -- Slide 2: Agenda --
    pcontent(prs, 'Agenda', [
        '1. Kontekst: dlaczego Fabric na polu walki?',
        '2. Scenariusz: OPERATION IRONSHIELD',
        '3. Faza 1: Real-Time Intelligence – świadomość sytuacyjna',
        '4. Faza 2: Data Agents – wywiad konwersacyjny',
        '5. Faza 3: Operations Agents – autonomiczne reagowanie',
        '6. Architektura techniczna',
        '7. Live Demo',
        '8. Wartość i kolejne kroki',
    ])

    # -- Slide 3: Context --
    pcontent(prs, 'Dlaczego Fabric na polu walki?', [
        'Pole walki generuje MILIONY zdarzeń/sekundę (drony, radary, GPS, wearable)',
        'Przewaga informacyjna = przewaga taktyczna (skrócenie pętli OODA)',
        'Potrzeba: jednolita platforma łącząca ingestion → analytics → AI → akcja',
        'Microsoft Fabric = zunifikowana platforma danych + AI agents',
        'Fabric RTI: sub-sekundowa analityka na miliardach zdarzeń',
        'Data Agents: dowódca pyta w języku naturalnym, agent odpowiada danymi',
        'Operations Agents: autonomiczny cykl observe → analyze → decide → act',
    ])

    # -- Slide 4: Scenario overview --
    p2col(prs, 'OPERATION IRONSHIELD – Scenariusz',
        'BLUE FORCE (siły sojusznicze)',
        [
            '24 pojazdy opancerzone (BWP Borsuk/Rosomak)',
            '12 dronów rozpoznawczych (FlyEye/Warmate)',
            '120 żołnierzy piechoty z wearable sensors',
            '4 radary artyleryjskie',
            '8 haubic samobieżnych Krab',
            '10 pojazdów logistycznych',
        ],
        'RED FORCE (symulowany przeciwnik)',
        [
            'Kolumna pancerna: 14 pojazdów',
            'Drony rekonesansowe',
            'Ataki artyleryjskie',
            'Grupy dywersyjne',
            'Walka elektroniczna (EW)',
            'Zachowanie parametryzowane w symulatorze',
        ]
    )

    # -- Slide 5: Architecture --
    pcontent(prs, 'Architektura techniczna', [
        'ŹRÓDŁA: Drony • Radary • Pojazdy • Żołnierze • Meteo • SIGINT • Logistyka',
        '↓ Eventstream (7 strumieni, Event Hub, MQTT/Kafka) ↓',
        'Eventhouse (KQL Database) – sub-sekundowe zapytania, mld rekordów',
        'Real-Time Dashboard – mapa operacyjna, KPI, alerty live',
        'Data Agents – NL → KQL, multi-source reasoning',
        'Operations Agents – observe→analyze→decide→act + human-in-the-loop',
        'Unity Catalog – governance, RBAC, audyt, lineage',
        'Power Automate – workflow: Teams alerts, akcje C2',
    ])

    # -- Slide 6: Phase 1 RTI --
    pcontent(prs, 'Faza 1: Real-Time Intelligence', [
        '🛰️  7 strumieni danych: drony, radary, pojazdy, żołnierze, meteo, SIGINT, logistyka',
        '⚡ Eventstream: ingestion milionów zdarzeń, transformacje inline',
        '🗄️  Eventhouse (KQL): sub-sekundowe zapytania analityczne',
        '🗺️  Real-Time Dashboard: mapa operacyjna z live pozycjami BLUE/RED',
        '🚨 Anomaly Detection: automatyczne wykrycie koncentracji sił wroga',
        '📢 Alert w Teams + na dashboardzie w < 5 sekund od wykrycia',
        '🔍 Drill-down: kliknięcie alertu → szczegóły KQL',
    ])

    # -- Slide 7: Phase 2 Data Agents --
    pcontent(prs, 'Faza 2: Data Agents – wywiad konwersacyjny', [
        '🗣️  Dowódca zadaje pytania w języku naturalnym',
        '🔄 Agent tłumaczy na KQL, odpytuje Eventhouse, syntetyzuje odpowiedź',
        '📊 Multi-source: korelacja radar + dron + SIGINT w jednej odpowiedzi',
        '🎖️  Intel Agent: "Ile wrogich kontaktów w sektorze Bravo? ETA?"',
        '📦 Supply Agent: "Stan amunicji 3. kompanii? ETA konwoju?"',
        '🩺 Medic Agent: "Czy są ranni żołnierze wymagający MEDEVAC?"',
        '⏱️  Czas odpowiedzi: sekundy zamiast minut ręcznej analizy',
    ])

    # -- Slide 8: Data Agent dialog example --
    pcontent(prs, 'Data Agent – przykład interakcji', [
        'DOWÓDCA: "Ile kontaktów wrogich w sektorze Bravo? Kierunek? Prędkość?"',
        '',
        'INTEL AGENT: "Wykryto 14 kontaktów hostile w sektorze Bravo."',
        '  → Śr. prędkość: 35 km/h, kierunek: NW (azymut 315°)',
        '  → Korelacja SIGINT: transmisje dowodzenia kolumną pancerną',
        '  → Dystans do BLUE: 12 km, ETA: 20 minut',
        '  → Rekomendacja: natychmiastowe alertowanie artylerii',
    ])

    # -- Slide 9: Phase 3 Operations Agents --
    pcontent(prs, 'Faza 3: Operations Agents – autonomiczne reagowanie', [
        '🔄 Pełny cykl: OBSERVE → ANALYZE → DECIDE → ACT',
        '👁️  OBSERVE: ciągłe monitorowanie KQL queries na Eventhouse',
        '🧠 ANALYZE: rozpoznanie sytuacji, kontekst (pogoda, siły, pozycje)',
        '💡 DECIDE: rekomendacja w Adaptive Card → Microsoft Teams',
        '✅ ACT: po zatwierdzeniu dowódcy → Power Automate workflow',
        '🛡️  Human-in-the-loop: AI rekomenduje, człowiek decyduje',
        '📋 Audyt: każda decyzja logowana w Unity Catalog',
    ])

    # -- Slide 10: Operations Agent - threat response --
    pcontent(prs, 'Operations Agent: Threat Response – pełny cykl', [
        '⚠️  OBSERVE: Radar wykrywa 14 hostile vehicles w sektorze Bravo',
        '🧠 ANALYZE: kolumna pancerna, heading NW, ETA 20 min, bateria Krab w zasięgu',
        '📱 DECIDE: Adaptive Card w Teams:',
        '   "Rekomendacja: ogień zaporowy artylerii na współrzędne 51.234°N, 21.456°E"',
        '   [✅ ZATWIERDŹ]   [❌ ODRZUĆ]   [✏️ MODYFIKUJ]',
        '✅ ACT: współrzędne → system ogniowy, dron BDA, log do Unity Catalog',
        '📊 VERIFY: dron potwierdza efekty na dashboardzie',
    ])

    # -- Slide 11: Operations Agent - logistics --
    pcontent(prs, 'Operations Agent: Logistics Optimizer', [
        '📉 DETECT: 3. kompania – stan amunicji 25% (< próg 30%)',
        '🗺️  ANALYZE: konwój KZ-04 w trasie, ETA 45 min, pozycje wroga na trasie',
        '🔀 RECOMMEND: zmiana trasy konwoju (omijając sektor Bravo)',
        '   + priorytet ładunku: amunicja > paliwo',
        '📱 Teams → oficer logistyki zatwierdza jednym kliknięciem',
        '🚛 ACT: aktualizacja routingu, nowe ETA: 55 min (bezpieczna trasa)',
        '📋 Unity Catalog: pełny audyt decyzji logistycznej',
    ])

    # -- Slide 12: Timeline --
    pcontent(prs, 'Oś czasu scenariusza (demo ~20 min)', [
        'T+0:00  Start sensorów, Eventstream live',
        'T+2:00  Dashboard: normalna aktywność, patrol',
        'T+5:00  ⚠️ Anomalia! Kontakty w sektorze Bravo',
        'T+6:00  Data Agent: "Ile kontaktów? Klasyfikacja?"',
        'T+9:00  Operations Agent: rekomendacja ognia artylerii',
        'T+10:00 Dowódca zatwierdza w Teams → workflow',
        'T+15:00 BDA: dron weryfikuje, Data Agent podsumowuje',
        'T+18:00 Logistics Agent: alert niski stan amunicji → reroute',
    ])

    # -- Slide 13: Value --
    p2col(prs, 'Wartość Microsoft Fabric na polu walki',
        'Zdolności',
        [
            'Pętla OODA: z minut do sekund',
            'Jedna platforma: dane + AI + akcja',
            'Human-in-the-loop: kontrola dowódcy',
            'Governance: audyt każdej decyzji',
            'Skalowalność: brygada → dywizja',
            'Interoperacyjność: REST API, Power Platform',
        ],
        'Porównanie z tradycyjnym C2',
        [
            'Ręczna analiza → AI-assisted',
            'Silosowe systemy → unified platform',
            'Reaktywne → proaktywne (agents)',
            'Raportowanie → real-time awareness',
            'Brak audytu → pełna traceability',
            'Pojedynczy użytkownik → multi-agent collaboration',
        ]
    )

    # -- Slide 14: Requirements --
    pcontent(prs, 'Wymagania techniczne demo', [
        'Microsoft Fabric: Capacity F2+, Copilot/AI enabled',
        'Workspace: dedykowany „IRONSHIELD"',
        'Eventstream: 7 strumieni + Azure Event Hub',
        'Eventhouse: 1 KQL Database (7 tabel)',
        'Data Agents: 3 (Intel, Supply, Medic)',
        'Operations Agents: 4 (Threat, Logistics, MEDEVAC, EW)',
        'Power Automate: 2 flows + Microsoft Teams kanał',
        'Generator danych: Python notebook (parametryzowany scenariusz)',
    ])

    # -- Slide 15: Next steps --
    pcontent(prs, 'Kolejne kroki', [
        '1. Provisioning: Fabric workspace + Event Hub + Eventhouse',
        '2. Implementacja generatora danych (Python notebook)',
        '3. Konfiguracja Eventstream (7 strumieni + transformacje)',
        '4. Build Real-Time Dashboard (7 paneli, mapa operacyjna)',
        '5. Konfiguracja Data Agents (3 agenty, sources, prompts)',
        '6. Konfiguracja Operations Agents (4 agenty, goals, actions)',
        '7. Power Automate flows + Teams kanał',
        '8. Rehearsal i optymalizacja scenariusza',
    ])

    # -- Slide 16: Thank you --
    ptitle(prs, 'OPERATION IRONSHIELD', 'Gotowi do demonstracji?\nPytania i dyskusja')

    path = os.path.join(OUT, 'IRONSHIELD_Scenario_Demo.pptx')
    prs.save(path)
    print(f"Created: {path}")
    return path


# ================================================================
# MAIN
# ================================================================

if __name__ == '__main__':
    print("Generating OPERATION IRONSHIELD documents...")
    create_scenario_docx()
    create_scenario_pptx()
    print("\nAll documents generated successfully!")
